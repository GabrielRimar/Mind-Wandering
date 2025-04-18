import subprocess
import os
from pptx import Presentation

class PresentationHandler:
    def __init__(self, file_path):
        self.presentation_path = file_path
        self.presentation = Presentation(file_path)

    def open_presentation(self):
        if os.path.exists(self.presentation_path):
            try:
                subprocess.run(["open", self.presentation_path])
                print("Presentation opened (or is opening)...")
            except Exception as e:
                print(f"Failed to open presentation: {e}")
        else:
            print("File not found. Please check the path.") 
    
    def get_number_of_slides(self):
        return len(self.presentation.slides)
    
    def get_number_of_words_per_slide(self):
        word_count = []
        for slide in self.presentation.slides:
            words = 0
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    words += len(text.split())
            word_count.append(words)
        return word_count
    def close_presentation(self):
        try:
            subprocess.run(["pkill", "-f", self.presentation_path])
            print("Presentation closed.")
        except Exception as e:
            print(f"Failed to close presentation: {e}")
    
    def save_data(self, file_path):
        import pandas as pd  # Only if not already at the top

        word_counts = self.get_number_of_words_per_slide()
        num_slides = self.get_number_of_slides()

        # Create a DataFrame with slide numbers and word counts
        df = pd.DataFrame({
            "slide_number": list(range(1, num_slides + 1)),
            "word_count": word_counts
        })

        # Make sure the output folder exists
        os.makedirs(file_path, exist_ok=True)

        # Save the DataFrame as CSV
        df.to_csv(file_path, index=False)
        print(f"Word counts saved to {file_path}")
        return df

if __name__ == "__main__":
    presentation_path = "presentation.pptx"
    ppt = PresentationHandler(presentation_path)
    ppt.open_presentation()
    print("Number of slides:", ppt.get_number_of_slides())
    print("Words per slide:", ppt.get_number_of_words_per_slide())
